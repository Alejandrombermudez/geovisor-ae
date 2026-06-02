"""Analizador profundo del PPTX 'Bancolombia Intro.pptx'.

Extrae para cada slide:
  - todos los shapes con su tipo, posición (left/top en EMU y %), tamaño
  - texto con formato (font, size, color, bold, alignment)
  - imágenes con la referencia al archivo en slides_assets/
  - tablas/gráficos si los hay

Salida: slides_manifest.json + un .md humano-legible por slide.
"""

from __future__ import annotations

import json
import os
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

ROOT = Path(r"D:\AMAZONIA EMPRENDE\GeoAE\Bancolombia_interfaz\presentacion")
PPTX = ROOT / "Bancolombia Intro.pptx"
ASSETS_DIR = ROOT / "slides_assets"
OUT_JSON = ROOT / "slides_manifest.json"
OUT_MD = ROOT / "slides_breakdown.md"


def emu_to_px(emu: int) -> int:
    """1 inch = 914400 EMU = 96 px (CSS)."""
    return round(emu / 914400 * 96)


def color_to_hex(color) -> str | None:
    try:
        if color and color.rgb is not None:
            return f"#{str(color.rgb)}"
    except Exception:
        pass
    return None


def analyze_run(run) -> dict:
    f = run.font
    return {
        "text": run.text,
        "bold": f.bold,
        "italic": f.italic,
        "size_pt": float(f.size.pt) if f.size else None,
        "name": f.name,
        "color": color_to_hex(f.color),
    }


def analyze_paragraph(p) -> dict:
    return {
        "alignment": str(p.alignment) if p.alignment else None,
        "level": p.level,
        "runs": [analyze_run(r) for r in p.runs],
        "text": p.text,
    }


def analyze_shape(shape, slide_width_emu: int, slide_height_emu: int, rels_map: dict | None = None) -> dict:
    rels_map = rels_map or {}
    info: dict = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "shape_id": shape.shape_id,
    }
    # Geometry
    try:
        info["left_emu"] = int(shape.left) if shape.left is not None else None
        info["top_emu"] = int(shape.top) if shape.top is not None else None
        info["width_emu"] = int(shape.width) if shape.width is not None else None
        info["height_emu"] = int(shape.height) if shape.height is not None else None
        if info["left_emu"] is not None:
            info["left_pct"] = round(info["left_emu"] / slide_width_emu * 100, 2)
            info["top_pct"] = round(info["top_emu"] / slide_height_emu * 100, 2)
            info["width_pct"] = round(info["width_emu"] / slide_width_emu * 100, 2)
            info["height_pct"] = round(info["height_emu"] / slide_height_emu * 100, 2)
            info["left_px"] = emu_to_px(info["left_emu"])
            info["top_px"] = emu_to_px(info["top_emu"])
            info["width_px"] = emu_to_px(info["width_emu"])
            info["height_px"] = emu_to_px(info["height_emu"])
    except Exception as e:
        info["geometry_error"] = str(e)

    # Text
    if shape.has_text_frame:
        tf = shape.text_frame
        info["text"] = tf.text
        info["paragraphs"] = [analyze_paragraph(p) for p in tf.paragraphs]

    # Picture
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            info["picture"] = {
                "content_type": img.content_type,
                "ext": img.ext,
                "size_bytes": len(img.blob),
                "filename": getattr(img, "filename", None),
            }
        except Exception as e:
            info["picture_error"] = str(e)

    # Detect image fills on any shape (e.g. Freeforms with picture fill)
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        blips = shape._element.findall(".//a:blip", ns)
        rids: list[str] = []
        for blip in blips:
            rid = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if rid and rid in rels_map:
                rids.append(rels_map[rid])
        # de-dup keeping order
        seen = set()
        rids = [x for x in rids if not (x in seen or seen.add(x))]
        if rids:
            info["asset_file"] = f"slides_assets/{rids[0]}"
            if len(rids) > 1:
                info["asset_files_all"] = [f"slides_assets/{x}" for x in rids]
    except Exception:
        pass

    # Group
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["children"] = [
            analyze_shape(s, slide_width_emu, slide_height_emu, rels_map)
            for s in shape.shapes
        ]

    # Tables
    if shape.has_table:
        tbl = shape.table
        info["table"] = {
            "rows": len(tbl.rows),
            "cols": len(tbl.columns),
            "cells": [[cell.text for cell in row.cells] for row in tbl.rows],
        }

    return info


def main() -> None:
    prs = Presentation(str(PPTX))
    sw_emu = prs.slide_width
    sh_emu = prs.slide_height

    manifest = {
        "source": PPTX.name,
        "slide_width_emu": sw_emu,
        "slide_height_emu": sh_emu,
        "slide_width_px": emu_to_px(sw_emu),
        "slide_height_px": emu_to_px(sh_emu),
        "aspect_ratio": round(sw_emu / sh_emu, 4),
        "n_slides": len(prs.slides),
        "slides": [],
    }

    md_lines: list[str] = [
        "# Bancolombia Intro — Desglose de slides",
        "",
        f"- Slides: **{len(prs.slides)}**",
        f"- Dimensiones: {emu_to_px(sw_emu)} × {emu_to_px(sh_emu)} px (ratio {round(sw_emu/sh_emu, 4)})",
        f"- Assets en: `slides_assets/`",
        "",
    ]

    for i, slide in enumerate(prs.slides, start=1):
        slide_info: dict = {
            "index": i,
            "layout_name": slide.slide_layout.name if slide.slide_layout else None,
            "shapes": [],
        }
        # Resolve rels for media filenames
        rels_map = {}
        for rid, rel in slide.part.rels.items():
            target = rel.target_ref
            if target.startswith("../media/"):
                rels_map[rid] = target.split("/")[-1]

        for shape in slide.shapes:
            s_info = analyze_shape(shape, sw_emu, sh_emu, rels_map)
            slide_info["shapes"].append(s_info)

        manifest["slides"].append(slide_info)

        # MD section (recursive tree)
        md_lines.append(f"## Slide {i} - `{slide_info['layout_name']}`")
        md_lines.append("")

        def render(s: dict, depth: int) -> None:
            indent = "  " * depth
            head = f"{indent}- **{s.get('name')}** ({s.get('shape_type')})"
            if "left_pct" in s:
                head += (
                    f" pos {s['left_pct']}%,{s['top_pct']}% "
                    f"size {s['width_pct']}%x{s['height_pct']}%"
                )
            md_lines.append(head)
            if s.get("text"):
                for line in s["text"].splitlines():
                    if line.strip():
                        md_lines.append(f"{indent}  > {line}")
            if "asset_file" in s:
                md_lines.append(f"{indent}  IMG  `{s['asset_file']}`")
            if "table" in s:
                md_lines.append(
                    f"{indent}  TABLE {s['table']['rows']}x{s['table']['cols']}"
                )
            for child in s.get("children", []):
                render(child, depth + 1)

        for s in slide_info["shapes"]:
            render(s, 0)
        md_lines.append("")

    OUT_JSON.write_text(json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
    OUT_MD.write_text("\n".join(md_lines), encoding="utf-8")
    print(f"OK Manifest: {OUT_JSON}")
    print(f"OK Breakdown: {OUT_MD}")
    print(f"  {len(prs.slides)} slides, {len(list(ASSETS_DIR.iterdir()))} assets")


if __name__ == "__main__":
    main()
